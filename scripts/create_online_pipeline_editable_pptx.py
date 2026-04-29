#!/usr/bin/env python3
"""
Create an editable PPTX slide for the clean online-policy pipeline.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "figures" / "architecture" / "online_gru_adaptive_pipeline_clean_editable.pptx"


def add_box(slide, x, y, w, h, title, subtitle, fill, border):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border
    shp.line.width = Pt(1.2)

    tf = shp.text_frame
    tf.clear()
    tf.word_wrap = True

    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = title
    r1.font.bold = True
    r1.font.size = Pt(12)
    r1.font.color.rgb = RGBColor(0x1F, 0x2A, 0x38)

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = subtitle
    r2.font.size = Pt(9)
    r2.font.color.rgb = RGBColor(0x3A, 0x3A, 0x3A)
    return shp


def add_line(slide, x1, y1, x2, y2, color=RGBColor(0x4D, 0x5A, 0x6A), width=Pt(1.2)):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    ln.line.color.rgb = color
    ln.line.width = width
    return ln


def add_label(slide, text, x, y, w=Inches(1.4), h=Inches(0.25), size=9):
    tx = slide.shapes.add_textbox(x, y, w, h)
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = RGBColor(0x2B, 0x3A, 0x4A)
    return tx


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.5))
    tf = title.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "Online Policy Adaptation Loop (Clean Editable)"
    r.font.bold = True
    r.font.size = Pt(28)
    r.font.color.rgb = RGBColor(0x00, 0x66, 0xCC)

    # Divider
    div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0.62), Inches(12.5), Inches(0.03))
    div.fill.solid()
    div.fill.fore_color.rgb = RGBColor(0x00, 0x66, 0xCC)
    div.line.fill.background()

    # Top row
    b1 = add_box(slide, Inches(0.3), Inches(0.75), Inches(2.1), Inches(0.72), "Shared Features", "Same cleaned feature pipeline", RGBColor(0xDF, 0xEC, 0xFF), RGBColor(0x9D, 0xB6, 0xD8))
    b2 = add_box(slide, Inches(2.55), Inches(0.75), Inches(2.1), Inches(0.72), "Window Timeline", "Build rolling windows + schedule", RGBColor(0xDF, 0xEC, 0xFF), RGBColor(0x9D, 0xB6, 0xD8))
    b3 = add_box(slide, Inches(4.8), Inches(0.75), Inches(2.1), Inches(0.72), "Per-Step History Slice", "Expanding/fixed lookback to train_end", RGBColor(0xDF, 0xEC, 0xFF), RGBColor(0x9D, 0xB6, 0xD8))
    b4 = add_box(slide, Inches(7.05), Inches(0.75), Inches(2.1), Inches(0.72), "Warm-Start Update", "Fine-tune current GRU for epochs_step", RGBColor(0xDD, 0xF5, 0xE7), RGBColor(0x9C, 0xCB, 0xB1))
    b5 = add_box(slide, Inches(9.3), Inches(0.75), Inches(2.6), Inches(0.72), "Update Policy", "Cadence + gate rules", RGBColor(0xFF, 0xF3, 0xD6), RGBColor(0xD8, 0xBE, 0x83))
    add_line(slide, Inches(2.4), Inches(1.11), Inches(2.55), Inches(1.11))
    add_line(slide, Inches(4.65), Inches(1.11), Inches(4.8), Inches(1.11))
    add_line(slide, Inches(6.9), Inches(1.11), Inches(7.05), Inches(1.11))
    add_line(slide, Inches(9.15), Inches(1.11), Inches(9.3), Inches(1.11))

    # Loop container
    loop = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.35), Inches(1.65), Inches(12.2), Inches(4.3))
    loop.fill.solid()
    loop.fill.fore_color.rgb = RGBColor(0xF9, 0xFB, 0xFD)
    loop.line.color.rgb = RGBColor(0xA9, 0xB5, 0xC2)
    loop.line.width = Pt(1.2)
    loop.text_frame.clear()

    add_label(slide, "Per scheduled decision date (t)", Inches(5.1), Inches(1.72), Inches(2.3), Inches(0.25), 10)

    # Inner boxes
    u1 = add_box(slide, Inches(0.8), Inches(2.1), Inches(2.0), Inches(0.8), "Is this an update step?", "If no, keep previous params", RGBColor(0xFF, 0xF3, 0xD6), RGBColor(0xD8, 0xBE, 0x83))
    u2 = add_box(slide, Inches(3.1), Inches(2.1), Inches(1.8), Inches(0.8), "Train/Val Split", "Chronological split", RGBColor(0xDF, 0xEC, 0xFF), RGBColor(0x9D, 0xB6, 0xD8))
    u3 = add_box(slide, Inches(5.2), Inches(2.1), Inches(1.95), Inches(0.8), "Update quality metrics", "val-loss delta + relative change", RGBColor(0xDF, 0xEC, 0xFF), RGBColor(0x9D, 0xB6, 0xD8))
    u4 = add_box(slide, Inches(7.45), Inches(2.1), Inches(1.9), Inches(0.8), "Accept/revert update?", "Threshold check (optional)", RGBColor(0xEF, 0xEF, 0xF6), RGBColor(0xB9, 0xB9, 0xC9))
    u5 = add_box(slide, Inches(4.6), Inches(3.5), Inches(2.25), Inches(0.8), "Infer weights for date t", "Action inference", RGBColor(0xDD, 0xF5, 0xE7), RGBColor(0x9C, 0xCB, 0xB1))
    u6 = add_box(slide, Inches(7.2), Inches(3.5), Inches(2.4), Inches(0.8), "Record realized outcomes", "Path accounting: return/costs/flags", RGBColor(0xF3, 0xE8, 0xFF), RGBColor(0xC6, 0xB0, 0xDE))

    # Branch and join lines
    add_line(slide, Inches(2.8), Inches(2.5), Inches(3.1), Inches(2.5))
    add_label(slide, "yes-update", Inches(2.82), Inches(2.26), Inches(1.2))
    add_line(slide, Inches(4.9), Inches(2.5), Inches(5.2), Inches(2.5))
    add_line(slide, Inches(7.15), Inches(2.5), Inches(7.45), Inches(2.5))
    add_line(slide, Inches(9.35), Inches(2.9), Inches(5.72), Inches(3.5))  # accepted/reverted params to inference
    add_label(slide, "accepted/reverted params", Inches(7.5), Inches(3.05), Inches(2.0))

    add_line(slide, Inches(1.8), Inches(2.9), Inches(1.8), Inches(3.9))
    add_line(slide, Inches(1.8), Inches(3.9), Inches(4.6), Inches(3.9))
    add_label(slide, "no-update", Inches(2.2), Inches(3.65), Inches(1.0))

    add_line(slide, Inches(6.85), Inches(3.9), Inches(7.2), Inches(3.9))

    # Single t->t+1 loop line
    add_line(slide, Inches(8.4), Inches(3.5), Inches(0.95), Inches(2.9))
    add_label(slide, "next decision date (t+1)", Inches(3.2), Inches(3.1), Inches(2.3))

    # Top hint from policy to loop
    add_line(slide, Inches(10.6), Inches(1.47), Inches(1.8), Inches(2.1))
    add_label(slide, "configured cadence + gate rules", Inches(7.9), Inches(1.6), Inches(2.7))

    # Output
    out = add_box(
        slide, Inches(3.1), Inches(6.15), Inches(6.2), Inches(0.8),
        "Online Outputs", "online path + schedule log + diagnostics + update-benefit summary",
        RGBColor(0xEE, 0xDF, 0xFF), RGBColor(0xC8, 0xAF, 0xDE)
    )
    add_line(slide, Inches(8.4), Inches(4.3), Inches(6.2), Inches(6.15))

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
