#!/usr/bin/env python3
"""
Append a slide to the architecture deck that reflects the latest actual online run.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
PPTX_PATH = ROOT / "figures" / "architecture" / "model_architecture_editable.pptx"
IMG_PATH = ROOT.parent / ".cursor" / "projects" / "Users-oliu-STAT-4830-OSO" / "assets" / "Screenshot_2026-04-23_at_2.01.31_PM-bdbed3a2-f4c3-4d57-9443-2b78b3b537aa.png"


def add_textbox(slide, text, left, top, width, height, size=14, bold=False, color=RGBColor(0, 0, 0), align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def main() -> None:
    if not PPTX_PATH.exists():
        raise FileNotFoundError(f"Missing PPTX: {PPTX_PATH}")
    if not IMG_PATH.exists():
        raise FileNotFoundError(f"Missing image: {IMG_PATH}")

    prs = Presentation(str(PPTX_PATH))
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # Title
    add_textbox(
        slide,
        "Online Pipeline — Actual Run Behavior (Latest)",
        Inches(0.6), Inches(0.25), Inches(15.0), Inches(0.6),
        size=30, bold=True, color=RGBColor(0x00, 0x66, 0xCC),
    )
    line = slide.shapes.add_shape(1, Inches(0.6), Inches(0.95), Inches(14.8), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0x00, 0x66, 0xCC)
    line.line.fill.background()

    # Main diagram image
    slide.shapes.add_picture(str(IMG_PATH), Inches(0.6), Inches(1.2), width=Inches(14.8), height=Inches(6.2))

    # Actual-run callout (no hard threshold acceptance gate)
    box = slide.shapes.add_shape(1, Inches(0.9), Inches(7.55), Inches(14.2), Inches(1.1))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xF2, 0xF8, 0xFF)
    box.line.color.rgb = RGBColor(0x99, 0xBB, 0xDD)

    add_textbox(
        slide,
        "Actual run notes: weekly updates (update_freq=W), epochs_step=2, warm-start fine-tuning, and no explicit hard "
        "accept/revert threshold confirmation gate in this run (updates applied on scheduled dates).",
        Inches(1.05), Inches(7.73), Inches(13.9), Inches(0.75),
        size=16, color=RGBColor(0x22, 0x22, 0x22),
    )

    # Highlight threshold-check caveat
    warn = slide.shapes.add_shape(1, Inches(10.55), Inches(0.18), Inches(4.8), Inches(0.45))
    warn.fill.solid()
    warn.fill.fore_color.rgb = RGBColor(0xFF, 0xEE, 0xCC)
    warn.line.color.rgb = RGBColor(0xDD, 0xAA, 0x66)
    add_textbox(
        slide,
        "Threshold-check confirm not used in latest run",
        Inches(10.7), Inches(0.24), Inches(4.5), Inches(0.30),
        size=12, bold=True, color=RGBColor(0x99, 0x44, 0x00), align=PP_ALIGN.CENTER,
    )

    prs.save(str(PPTX_PATH))
    print(f"Updated deck: {PPTX_PATH}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
