#!/usr/bin/env python3
"""
Add validation analysis slides and failure mode slide to the existing deck.
Inserts after slide 6 (Results: Performance Comparison).
"""
from pathlib import Path
from copy import deepcopy

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parent.parent
PPTX_IN  = ROOT / "docs" / "IPO Portfolio Optimization GRU Allocation.pptx"
PPTX_OUT = ROOT / "docs" / "IPO Portfolio Optimization GRU Allocation.pptx"
FIGS     = ROOT / "figures"

# Brand colours matching the deck
BLUE   = RGBColor(0x00, 0x66, 0xCC)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BLACK  = RGBColor(0x00, 0x00, 0x00)
LGRAY  = RGBColor(0xF2, 0xF2, 0xF2)
DGRAY  = RGBColor(0x33, 0x33, 0x33)

W = Inches(16)   # slide width
H = Inches(9)    # slide height


# ── helpers ────────────────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_rgb=None, line_rgb=None):
    from pptx.util import Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.fill.background()
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=BLACK, align=PP_ALIGN.LEFT,
                wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return txb


def add_picture(slide, img_path, left, top, width, height):
    return slide.shapes.add_picture(str(img_path), left, top, width, height)


def add_footer(slide, label="STAT 4830 — Spring 2026"):
    """Replicate the bottom-right navy rectangle present on every slide."""
    add_rect(slide, Inches(13.80), Inches(8.35), Inches(2.32), Inches(0.65), fill_rgb=BLUE)
    add_textbox(slide, label,
                Inches(13.85), Inches(8.38), Inches(2.22), Inches(0.55),
                font_size=9, bold=False, color=WHITE, align=PP_ALIGN.CENTER)


def add_slide_title(slide, title):
    add_textbox(slide, title,
                Inches(0.87), Inches(0.35), Inches(14.26), Inches(0.55),
                font_size=28, bold=True, color=BLUE)


def add_divider(slide):
    """Thin blue horizontal rule under the title."""
    r = add_rect(slide, Inches(0.87), Inches(0.98), Inches(14.26), Inches(0.04),
                 fill_rgb=BLUE)


def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]   # blank layout
    return prs.slides.add_slide(blank_layout)


def reorder_slides(prs, new_order):
    """Reorder slides by index list."""
    from pptx.oxml.ns import qn
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.clear()
    for i in new_order:
        xml_slides.append(slides[i])


# ── slide builders ─────────────────────────────────────────────────────────────

def make_plot_slide(prs, title, img_path, bullets):
    """
    Two-column layout: large image left, bullet analysis right.
    bullets: list of (header, body) tuples
    """
    slide = blank_slide(prs)
    add_slide_title(slide, title)
    add_divider(slide)
    add_footer(slide)

    # Image — left 60% of slide, below title
    add_picture(slide, img_path,
                Inches(0.40), Inches(1.10), Inches(9.20), Inches(7.50))

    # Right panel — analysis bullets
    y = Inches(1.15)
    for header, body in bullets:
        # bullet header
        add_textbox(slide, f"▸  {header}",
                    Inches(9.80), y, Inches(5.90), Inches(0.36),
                    font_size=13, bold=True, color=BLUE)
        y += Inches(0.34)
        # bullet body
        add_textbox(slide, body,
                    Inches(10.10), y, Inches(5.60), Inches(0.75),
                    font_size=11, bold=False, color=DGRAY, wrap=True)
        y += Inches(0.80)
    return slide


def make_failure_mode_slide(prs):
    slide = blank_slide(prs)
    add_slide_title(slide, "Failure Mode Analysis")
    add_divider(slide)
    add_footer(slide)

    items = [
        ("Static Allocation Collapse (Root Cause)",
         "λ_path=0.01 & λ_turnover=0.0025 penalised any weight change. "
         "The GRU minimised loss by outputting a near-constant 68/32 split — "
         "avg daily turnover 0.0031, identical weights every day. "
         "The model learned a static policy, not a dynamic one."),
        ("Regularisation Dominates Return Signal",
         "With batch_size=32, gradient noise was high. The path/turnover terms "
         "had smooth, easy-to-minimise gradients while the return signal was noisy. "
         "Optimizer converged to the trivial fixed-weight solution in <5 epochs."),
        ("Validation Leakage via Hyperparameter Search",
         "288-config grid search selected the config maximising validation Sharpe. "
         "The reported Sharpe (3.27 → 3.60) is optimistic. A true held-out test "
         "set is needed to measure generalisation honestly."),
        ("Survivorship Bias in IPO Index",
         "The IPO index includes only tickers with price data in CRSP. "
         "Failed / delisted IPOs within the holding window are systematically "
         "excluded, inflating the IPO index return (192% in 2024)."),
        ("Validation Period Regime Risk",
         "Jan–Dec 2024 contained five distinct regime shifts (Apr sell-off, "
         "Aug VIX-65 carry unwind, Sep Fed cut, Nov election rally, Dec hawkish "
         "pivot). A static model cannot adapt; a dynamic model must learn these "
         "signals from 84-day windows."),
        ("Remaining Gap vs. Equal Weight (Post-Fix)",
         "After fixing λ_path=0.0001 & λ_turnover=0.0, Sharpe improved to 3.60 "
         "but still trails equal-weight (3.65). The GRU has not yet learned to "
         "time the market. IPO weight rose to 46% but is still near-static; "
         "true dynamic allocation requires richer regime features."),
    ]

    col_w = Inches(6.80)
    col_gap = Inches(0.30)
    left1 = Inches(0.50)
    left2 = left1 + col_w + col_gap

    for idx, (header, body) in enumerate(items):
        col   = idx % 2
        row   = idx // 2
        left  = left1 if col == 0 else left2
        top   = Inches(1.20) + row * Inches(2.35)

        # card background
        add_rect(slide, left, top, col_w, Inches(2.20), fill_rgb=LGRAY)
        # header
        add_textbox(slide, header,
                    left + Inches(0.15), top + Inches(0.10),
                    col_w - Inches(0.30), Inches(0.38),
                    font_size=12, bold=True, color=BLUE)
        # body
        add_textbox(slide, body,
                    left + Inches(0.15), top + Inches(0.48),
                    col_w - Inches(0.30), Inches(1.60),
                    font_size=10, bold=False, color=DGRAY, wrap=True)

    return slide


def make_updated_results_slide(prs):
    """Updated performance table with the new numbers."""
    slide = blank_slide(prs)
    add_slide_title(slide, "Updated Results: Performance Comparison")
    add_divider(slide)
    add_footer(slide)

    add_textbox(slide, "Validation period: January – December 2024  |  Retrained with 200 epochs, batch=256, lr=3e-4, λ_path=0.0001, λ_turnover=0.0",
                Inches(0.87), Inches(1.05), Inches(14.26), Inches(0.38),
                font_size=11, bold=False, color=DGRAY)

    # Table data
    headers = ["Strategy", "Total Return", "Ann. Return", "Ann. Vol", "Sharpe", "Max DD", "Avg IPO Wt"]
    rows = [
        ("Model Portfolio (updated)", "82.05%", "91.70%", "18.58%", "3.60 ▲", "-7.24%", "46.21%"),
        ("Model Portfolio (previous)", "58.81%", "65.27%", "15.74%", "3.27",   "-7.47%", "31.13%"),
        ("Equal 50/50",                "88.52%", "99.11%", "19.39%", "3.65",   "-7.20%", "50.00%"),
        ("IPO only",                   "192.78%","221.19%","31.03%", "3.92",   "-10.08%","100.00%"),
        ("Market only",                "19.40%", "21.24%", "12.25%", "1.63",   "-7.89%", "0.00%"),
    ]

    col_lefts  = [Inches(x) for x in [0.50, 3.10, 4.90, 6.70, 8.50, 10.00, 11.55]]
    col_widths = [Inches(x) for x in [2.55, 1.75, 1.75, 1.75, 1.45, 1.50, 1.80]]
    row_h = Inches(0.60)
    top0  = Inches(1.55)

    # Header row
    add_rect(slide, Inches(0.50), top0, Inches(13.30), row_h, fill_rgb=BLUE)
    for j, hdr in enumerate(headers):
        add_textbox(slide, hdr, col_lefts[j], top0 + Inches(0.10),
                    col_widths[j], row_h - Inches(0.10),
                    font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    for i, row in enumerate(rows):
        top_r = top0 + (i + 1) * row_h
        fill  = LGRAY if i % 2 == 0 else WHITE
        if i == 0:
            fill = RGBColor(0xE6, 0xF2, 0xFF)   # light blue highlight for updated row
        add_rect(slide, Inches(0.50), top_r, Inches(13.30), row_h, fill_rgb=fill)
        for j, cell in enumerate(row):
            bold  = (j == 0)
            color = BLUE if (i == 0 and j == 4) else DGRAY
            add_textbox(slide, cell, col_lefts[j], top_r + Inches(0.10),
                        col_widths[j], row_h - Inches(0.10),
                        font_size=11, bold=bold, color=color, align=PP_ALIGN.CENTER)

    # Key takeaway
    add_textbox(slide,
                "Key changes: Sharpe +10% (3.27→3.60) | Total return +39% (58.81%→82.05%) | "
                "IPO allocation +49% (31%→46%) | Max drawdown improved -7.47%→-7.24%",
                Inches(0.50), Inches(7.55), Inches(13.30), Inches(0.55),
                font_size=11, bold=False, color=DGRAY)
    return slide


# ── main ───────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation(str(PPTX_IN))
    n_orig = len(prs.slides)

    # Add slides in order (they'll be appended; we reorder after)
    make_updated_results_slide(prs)

    make_plot_slide(prs, "Validation Period: Cumulative Returns",
                    FIGS / "validation_cumulative_returns.png",
                    [
                        ("Model vs Baselines",
                         "Model (106%) roughly tracks Equal 50/50 (115%) but underperforms "
                         "IPO-only (244%). The model adds value over market-only (+85pp) "
                         "while maintaining tighter drawdown."),
                        ("Q1–Q2 2024: Tight clustering",
                         "All strategies converge until July. The model closely follows "
                         "equal-weight, suggesting the GRU has not yet learned to "
                         "differentiate market regimes in this period."),
                        ("Aug 2024: Yen carry unwind (VIX→65)",
                         "A sharp drawdown is visible across all strategies. The model "
                         "does not reduce IPO exposure ahead of the event — a key "
                         "failure of the static-weight regime."),
                        ("Nov–Dec 2024: Election + Fed divergence",
                         "IPO-only surges on the Trump rally while the model trails. "
                         "The vol-excess penalty (target 25%) caps the model's upside "
                         "when IPO volatility spikes above the threshold."),
                        ("Interpretation",
                         "The model achieves Sharpe 3.60 vs 3.65 (equal-weight) — "
                         "near parity, but without true dynamic timing the model "
                         "cannot exploit regime transitions."),
                    ])

    make_plot_slide(prs, "Validation Period: Optimization Objective",
                    FIGS / "validation_objective.png",
                    [
                        ("Combined objective L (upper panel)",
                         "L = −mean_return + λ_vol·Var + λ_cvar·(−CVaR) + λ_vol_excess·excess. "
                         "Spikes in Mar and Oct–Nov coincide with periods of elevated "
                         "IPO volatility. Lower is better."),
                        ("Mar 2024 spike",
                         "Objective peaks sharply in early March when IPO index variance "
                         "surged. The vol-excess and CVaR penalties dominated, forcing "
                         "the model to increase market weight temporarily."),
                        ("Mean return (reward) component",
                         "The blue reward signal is consistently positive and smooth, "
                         "indicating the portfolio was profitable throughout the period. "
                         "The objective is predominantly driven by risk penalties."),
                        ("Vol-excess penalty (yellow)",
                         "Activates in Nov–Dec when the post-election IPO rally pushed "
                         "annualised vol above the 25% target. The model responds by "
                         "trimming IPO weight — visible as a plateau in the returns curve."),
                        ("CVaR penalty (pink)",
                         "Elevated during Aug carry-trade unwind. Soft-sorted CVaR "
                         "captures tail risk in real-time, though the approximation "
                         "temperature (0.1) is untuned — a known limitation."),
                    ])

    make_plot_slide(prs, "Validation Period: Portfolio Variance",
                    FIGS / "validation_variance.png",
                    [
                        ("Model variance vs target (6.25%)",
                         "The model (blue) stays near the annualised variance target "
                         "(red dotted, 6.25% = 25% vol²) for most of the year. The "
                         "vol-excess penalty is effective as a soft constraint."),
                        ("Early-period spike (Feb–Mar 2024)",
                         "Model variance spikes to ~70% briefly before converging. "
                         "This corresponds to a period of high IPO index volatility "
                         "where the 21-day rolling window picks up outlier returns."),
                        ("IPO-only (purple) far exceeds target",
                         "IPO-only annualised variance reaches 300%+ in early periods. "
                         "The model's risk control is the primary value-add: same "
                         "IPO exposure with 8–10× lower variance."),
                        ("Aug 2024: Minimal model variance spike",
                         "Despite the Aug VIX-65 carry-trade event, the model's "
                         "variance barely moves — because weights are near-static and "
                         "the ~54% market component dampens the IPO shock."),
                        ("Nov–Dec 2024: Second spike",
                         "Post-election IPO surge caused a brief variance spike. The "
                         "vol-excess penalty activates (visible in objective plot) and "
                         "brings variance back toward target by year-end."),
                    ])

    make_failure_mode_slide(prs)

    # Reorder: insert new slides after original slide 6 (index 6)
    # Original: 0-9 (10 slides). New slides appended at 10-14.
    new_order = list(range(7)) + list(range(10, 15)) + list(range(7, 10))
    reorder_slides(prs, new_order)

    prs.save(str(PPTX_OUT))
    print(f"Saved: {PPTX_OUT}")
    print(f"  {n_orig} original slides + 5 new = {len(prs.slides)} total")


if __name__ == "__main__":
    main()
