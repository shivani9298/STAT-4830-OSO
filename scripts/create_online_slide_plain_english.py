#!/usr/bin/env python3
"""
Create an editable PPTX slide that explains the online pipeline in plain English.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "figures" / "architecture" / "online_slide_plain_english_editable.pptx"

BOX_EXPLANATIONS: list[tuple[str, list[str]]] = [
    (
        "Per scheduled decision date",
        [
            "For each day/week/month where a decision is allowed, run this mini-process once.",
        ],
    ),
    (
        "Update date? (If not, keep previous parameters)",
        [
            "Ask: Is today one of the dates where we retrain/update the model?",
            "If yes: try to improve the model with newest data.",
            "If no: reuse the model from last time.",
        ],
    ),
    (
        "yes branch",
        [
            "Used on update dates.",
            "Goes into model-check flow before new parameters are used.",
        ],
    ),
    (
        "Train/Val Split (Chronological split)",
        [
            "Use history up to today and split into train and validation.",
            "Train part fits/updates the model.",
            "Validation part checks if update is actually better.",
            "Chronological means past to recent order; no random shuffle.",
        ],
    ),
    (
        "Candidate Metrics (val-loss delta and relative change)",
        [
            "After update, measure whether validation improved.",
            "Delta = numeric change in validation loss.",
            "Relative change = percentage-style improvement/worsening.",
            "Question answered: Did this update help enough to trust?",
        ],
    ),
    (
        "Threshold Check (Cadence or confidence gate)",
        [
            "Decision gate: accept update only if it passes rule.",
            "Example: accept only if validation improves by at least X.",
            "If gate fails, revert and keep previous model.",
        ],
    ),
    (
        "no branch",
        [
            "Used when today is not an update date.",
            "Skips retraining checks and goes straight to allocation.",
        ],
    ),
    (
        "Action Inference (Predict allocation for current decision)",
        [
            "Use accepted model to output today's portfolio weights.",
            "Example output: 70% market, 30% IPO.",
        ],
    ),
    (
        "Path Accounting (Realized return, turnover, costs, flags)",
        [
            "After returns are known, record realized performance.",
            "Track return, turnover, transaction cost, update flag.",
            "Build the full backtest path row by row.",
        ],
    ),
]


def add_title(slide) -> None:
    title = slide.shapes.add_textbox(Inches(0.45), Inches(0.2), Inches(12.4), Inches(0.55))
    tf = title.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "Online Pipeline: Plain-English Box Definitions"
    r.font.bold = True
    r.font.size = Pt(28)
    r.font.color.rgb = RGBColor(0x00, 0x57, 0xA4)

    divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.45), Inches(0.72), Inches(12.2), Inches(0.03))
    divider.fill.solid()
    divider.fill.fore_color.rgb = RGBColor(0x00, 0x57, 0xA4)
    divider.line.fill.background()


def add_card(slide, x, y, w, h, heading: str, bullets: list[str]) -> None:
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xF7, 0xFA, 0xFD)
    card.line.color.rgb = RGBColor(0xB9, 0xC8, 0xD9)
    card.line.width = Pt(1.1)

    tf = card.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(6)

    hp = tf.paragraphs[0]
    hr = hp.add_run()
    hr.text = heading
    hr.font.bold = True
    hr.font.size = Pt(12)
    hr.font.color.rgb = RGBColor(0x1F, 0x2A, 0x38)

    for item in bullets:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(0x2F, 0x2F, 0x2F)
        p.space_before = Pt(1)
        p.space_after = Pt(1)


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title(slide)

    x_left = Inches(0.45)
    x_right = Inches(6.85)
    card_w = Inches(6.0)
    card_h = Inches(1.24)
    y0 = Inches(0.9)
    dy = Inches(1.28)

    left_items = BOX_EXPLANATIONS[:5]
    right_items = BOX_EXPLANATIONS[5:]

    for i, (title, bullets) in enumerate(left_items):
        add_card(slide, x_left, y0 + i * dy, card_w, card_h, title, bullets)

    for i, (title, bullets) in enumerate(right_items):
        add_card(slide, x_right, y0 + i * dy, card_w, card_h, title, bullets)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
